# -*- coding: utf-8 -*-
"""
Created on Mon Aug 21 08:42:39 2017

@author: BALASUBRAMANIAM
"""

from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

prs.save('test.pptx')